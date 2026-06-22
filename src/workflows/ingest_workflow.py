'''Ingest Workflow - Full document processing pipeline.'''

from __future__ import annotations

import asyncio
import logging
import time
import uuid
from pathlib import Path
from typing import Any

from loguru import logger

from src.core.document_object import DocumentObject, DocumentFormat
from src.core.geometric_element import ElementType, GeometricElement, make_element
from src.core.normalized_document import (
    BlockType, BoundingBox, NormalizedBlock, NormalizedDocument, NormalizedPage,
)
from src.engines.geometry.geometry_engine import GeometryEngine
from src.engines.recurrence.recurrence_engine import RecurrenceEngine
from src.engines.frequency.frequency_engine import FrequencyEngine
from src.engines.matrix.matrix_engine import MatrixEngine, TableMatrix
from src.engines.template.template_engine import TemplateEngine
from src.engines.graph.graph_engine import GraphEngine, GraphBuilder
from src.engines.semantic.semantic_engine import EmbeddingService, SemanticEngine
from src.engines.memory.hierarchical_memory import HierarchicalMemory
from src.engines.vector_db.faiss_store import FAISSStore as FaissStore
from src.ingestion.service import IngestionService
from src.normalization.cleaner import TextCleaner
from src.normalization.layout import LayoutDetector
from src.normalization.ocr import OCREngine


class IngestWorkflow:
    '''
    Complete document ingestion workflow.

    Phases:
    1. Load (multi-format)
    2. Parse + OCR + Normalize
    3. Convert to geometric elements
    4. Build all representations
    5. Index in vector store
    6. Store in hierarchical memory
    '''

    def __init__(
        self,
        llm_provider: str = 'openai',
        llm_model: str = 'gpt-4o-mini',
        llm_api_key: str = '',
        embedding_device: str = 'cpu',
        redis_url: str | None = None,
    ):
        self.ingestion = IngestionService()
        self.ocr = OCREngine()
        self.layout_detector = LayoutDetector()
        self.cleaner = TextCleaner()
        self.embedder = EmbeddingService(device=embedding_device)
        self.semantic = SemanticEngine(self.embedder)
        self.memory = HierarchicalMemory(redis_url or 'redis://localhost:6379/0')
        # Lazy init
        self.geometry = None
        self.recurrence = None
        self.frequency = None
        self.matrix = None
        self.template = None
        self.graph_builder = GraphBuilder()
        self.vector_store = FaissStore(dim=1024)
        # State
        self._elements: list[GeometricElement] = []
        self._tables: list[TableMatrix] = []
        self._graph: Any = None
        self._normalized: NormalizedDocument | None = None
        self._doc_id: str = ''
        self._filename: str = ''
        self._llm_provider = llm_provider
        self._llm_model = llm_model
        self._llm_api_key = llm_api_key

    async def initialize(self) -> None:
        '''Initialize async resources.'''
        await self.memory.connect()

    async def shutdown(self) -> None:
        '''Cleanup.'''
        await self.memory.close()

    async def ingest(
        self,
        source: str | Path | bytes | DocumentObject,
        filename: str = '',
        format: DocumentFormat | None = None,
    ) -> dict:
        '''Run full ingestion pipeline.'''
        timings = {}
        t_start = time.perf_counter()

        # Phase 1: Load
        t0 = time.perf_counter()
        if isinstance(source, DocumentObject):
            doc = source
        else:
            doc = await self.ingestion.ingest(source, filename, format)
        timings['phase1_load_s'] = round(time.perf_counter() - t0, 3)
        logger.info(f'Loaded: {doc.filename} ({doc.size_bytes} bytes)')

        self._doc_id = doc.doc_id
        self._filename = doc.filename

        # Phase 2: Normalize
        t0 = time.perf_counter()
        normalized = await self._normalize(doc)
        timings['phase2_normalize_s'] = round(time.perf_counter() - t0, 3)
        logger.info(f'Normalized: {normalized.total_pages} pages, {normalized.total_blocks} blocks')
        self._normalized = normalized

        # Phase 3: Convert to geometric elements
        t0 = time.perf_counter()
        elements = self._to_elements(normalized)
        timings['phase3_elements_s'] = round(time.perf_counter() - t0, 3)
        self._elements = elements
        logger.info(f'Created {len(elements)} geometric elements')

        # Phase 4-9: Build all representations
        t0 = time.perf_counter()
        # Geometry
        self.geometry = GeometryEngine()
        for page in normalized.pages:
            self.geometry.normalize_coordinates(page.page_number, page.width, page.height)
        self.geometry.add_many(elements)

        # Recurrence
        self.recurrence = RecurrenceEngine()
        self.recurrence.detect(elements)

        # Frequency weighting
        self.frequency = FrequencyEngine()
        self.frequency.assign_weights(elements)

        # Matrix extraction
        self.matrix = MatrixEngine()
        self._tables = self.matrix.find_tables(elements)

        # Template detection
        self.template = TemplateEngine()
        self.template.build(elements)

        # Graph
        self._graph = self.graph_builder.build(elements)
        timings['phase4_9_engines_s'] = round(time.perf_counter() - t0, 3)

        # Phase 9: Semantic indexing
        t0 = time.perf_counter()
        texts = [e.content for e in elements if e.content]
        metas = [{'element_id': e.element_id, 'page': e.page, 'section': e.section}
                 for e in elements if e.content]
        if texts:
            await self.semantic.process(texts, metas)
            embeddings = self.embedder.encode(texts)
            v_metas = [
                {'chunk_id': e.element_id, 'text': e.content, 'page': e.page,
                 'section': e.section, 'type': e.type.value, 'doc_id': e.doc_id}
                for e in elements if e.content
            ]
            await self.vector_store.upsert(embeddings, v_metas)
        timings['phase9_semantic_s'] = round(time.perf_counter() - t0, 3)
        logger.info('Semantic indexing complete')

        # Phase 11: Memory storage
        t0 = time.perf_counter()
        await self.memory.put(f'doc:{doc.doc_id}', {
            'filename': doc.filename,
            'elements': len(elements),
            'tables': len(self._tables),
            'templates': len(self.template.templates),
            'graph_nodes': self._graph.graph.number_of_nodes() if self._graph else 0,
            'graph_edges': self._graph.graph.number_of_edges() if self._graph else 0,
        })
        timings['phase11_memory_s'] = round(time.perf_counter() - t0, 3)

        # Total
        timings['total_s'] = round(time.perf_counter() - t_start, 3)

        return {
            'doc_id': doc.doc_id,
            'filename': doc.filename,
            'pages': normalized.total_pages,
            'blocks': normalized.total_blocks,
            'tables': len(self._tables),
            'templates': len(self.template.templates),
            'graph_nodes': self._graph.graph.number_of_nodes() if self._graph else 0,
            'graph_edges': self._graph.graph.number_of_edges() if self._graph else 0,
            'timings': timings,
        }

    async def _normalize(self, doc: DocumentObject) -> NormalizedDocument:
        '''Parse and normalize any document format.'''
        if doc.format == DocumentFormat.PDF:
            normalized = await self._normalize_pdf(doc)
        elif doc.format == DocumentFormat.DOCX:
            normalized = await self._normalize_docx(doc)
        elif doc.format == DocumentFormat.PPTX:
            normalized = await self._normalize_pptx(doc)
        elif doc.format == DocumentFormat.XLSX:
            normalized = await self._normalize_xlsx(doc)
        elif doc.format == DocumentFormat.IMAGE:
            normalized = await self._normalize_image(doc)
        else:
            normalized = await self._normalize_text(doc)
        return normalized

    async def _normalize_pdf(self, doc: DocumentObject) -> NormalizedDocument:
        import fitz
        normalized = NormalizedDocument(doc_id=doc.doc_id, filename=doc.filename)
        pdf = fitz.open(stream=doc.raw_bytes, filetype='pdf')
        for page_index in range(len(pdf)):
            page_obj = pdf[page_index]
            page = NormalizedPage(
                page_number=page_index + 1,
                width=page_obj.rect.width,
                height=page_obj.rect.height,
            )
            text_dict = page_obj.get_text('dict')
            for block in text_dict.get('blocks', []):
                if block.get('type') == 0:
                    text = self._reconstruct_text(block)
                    if not text.strip():
                        continue
                    bbox = block.get('bbox', (0, 0, 0, 0))
                    btype = self._classify_block(text, bbox, page.height)
                    page.blocks.append(NormalizedBlock(
                        type=btype,
                        text=self.cleaner.clean(text),
                        bbox=BoundingBox(*bbox[:4]),
                        page=page.page_number,
                    ))
                elif block.get('type') == 1:
                    bbox = block.get('bbox', (0, 0, 0, 0))
                    page.blocks.append(NormalizedBlock(
                        type=BlockType.FIGURE, text='',
                        bbox=BoundingBox(*bbox[:4]),
                        page=page.page_number,
                    ))
            # Tables
            self._extract_tables(page_obj, page)
            # OCR fallback
            if len(page.text.strip()) < 50:
                try:
                    pix = page_obj.get_pixmap(dpi=150)
                    page.page_image = pix.tobytes('png')
                    page.is_scanned = True
                    ocr_text = await self.ocr.recognize(page.page_image)
                    if ocr_text:
                        page.blocks.insert(0, NormalizedBlock(
                            type=BlockType.TEXT,
                            text=self.cleaner.clean(ocr_text),
                            page=page.page_number,
                        ))
                except Exception:
                    pass
            # Layout analysis
            page = self.layout_detector.analyze(page)
            normalized.pages.append(page)
        pdf.close()
        return normalized

    async def _normalize_docx(self, doc: DocumentObject) -> NormalizedDocument:
        import io
        from docx import Document as DocxDocument
        normalized = NormalizedDocument(doc_id=doc.doc_id, filename=doc.filename)
        d = DocxDocument(io.BytesIO(doc.raw_bytes))
        page = NormalizedPage(page_number=1, width=612.0, height=792.0)
        for para in d.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            btype = BlockType.TITLE if para.style.name.startswith('Heading 1') else \
                    BlockType.HEADING if para.style.name.startswith('Heading') else \
                    BlockType.TEXT
            page.blocks.append(NormalizedBlock(type=btype, text=text, page=1))
        for tbl in d.tables:
            rows = [[c.text.strip() for c in r.cells] for r in tbl.rows]
            page.blocks.append(NormalizedBlock(
                type=BlockType.TABLE,
                text=self._rows_to_md(rows),
                page=1,
            ))
        normalized.pages.append(page)
        return normalized

    async def _normalize_pptx(self, doc: DocumentObject) -> NormalizedDocument:
        import io
        from pptx import Presentation
        normalized = NormalizedDocument(doc_id=doc.doc_id, filename=doc.filename)
        pres = Presentation(io.BytesIO(doc.raw_bytes))
        for i, slide in enumerate(pres.slides, start=1):
            page = NormalizedPage(
                page_number=i,
                width=pres.slide_width,
                height=pres.slide_height,
            )
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = '\n'.join(p.text for p in shape.text_frame.paragraphs).strip()
                    if text:
                        btype = BlockType.TITLE if shape == slide.shapes.title else BlockType.TEXT
                        page.blocks.append(NormalizedBlock(type=btype, text=text, page=i))
                if shape.shape_type == 13:
                    page.blocks.append(NormalizedBlock(type=BlockType.FIGURE, text='', page=i))
            normalized.pages.append(page)
        return normalized

    async def _normalize_xlsx(self, doc: DocumentObject) -> NormalizedDocument:
        import io
        from openpyxl import load_workbook
        normalized = NormalizedDocument(doc_id=doc.doc_id, filename=doc.filename)
        wb = load_workbook(io.BytesIO(doc.raw_bytes), data_only=True)
        for i, sheet_name in enumerate(wb.sheetnames, start=1):
            ws = wb[sheet_name]
            page = NormalizedPage(page_number=i, width=612.0, height=792.0)
            rows = [[str(c.value) if c.value is not None else '' for c in row]
                    for row in ws.iter_rows()]
            if rows:
                page.blocks.append(NormalizedBlock(
                    type=BlockType.TABLE,
                    text=self._rows_to_md(rows),
                    page=i,
                ))
            normalized.pages.append(page)
        return normalized

    async def _normalize_image(self, doc: DocumentObject) -> NormalizedDocument:
        normalized = NormalizedDocument(doc_id=doc.doc_id, filename=doc.filename)
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(doc.raw_bytes))
        page = NormalizedPage(
            page_number=1, width=img.width, height=img.height,
            is_scanned=True,
        )
        ocr_text = await self.ocr.recognize(doc.raw_bytes)
        if ocr_text:
            page.blocks.append(NormalizedBlock(
                type=BlockType.TEXT, text=self.cleaner.clean(ocr_text), page=1,
            ))
        normalized.pages.append(page)
        return normalized

    async def _normalize_text(self, doc: DocumentObject) -> NormalizedDocument:
        normalized = NormalizedDocument(doc_id=doc.doc_id, filename=doc.filename)
        text = doc.raw_bytes.decode('utf-8', errors='ignore')
        page = NormalizedPage(page_number=1)
        for line in text.splitlines():
            line = line.strip()
            if line:
                page.blocks.append(NormalizedBlock(type=BlockType.TEXT, text=line, page=1))
        normalized.pages.append(page)
        return normalized

    def _to_elements(self, normalized: NormalizedDocument) -> list[GeometricElement]:
        type_map = {
            BlockType.TEXT: ElementType.TEXT, BlockType.TABLE: ElementType.TABLE,
            BlockType.FIGURE: ElementType.FIGURE, BlockType.EQUATION: ElementType.EQUATION,
            BlockType.HEADER: ElementType.HEADER, BlockType.FOOTER: ElementType.FOOTER,
            BlockType.CAPTION: ElementType.CAPTION, BlockType.LIST: ElementType.LIST_ITEM,
            BlockType.TITLE: ElementType.TITLE, BlockType.HEADING: ElementType.HEADING,
            BlockType.CODE: ElementType.CODE,
        }
        elements = []
        for page in normalized.pages:
            for block in page.blocks:
                bbox = None
                if block.bbox:
                    bbox = BoundingBox(block.bbox.x0, block.bbox.y0, block.bbox.x1, block.bbox.y1)
                etype = type_map.get(block.type, ElementType.TEXT)
                elements.append(make_element(
                    content=block.text, page=page.page_number,
                    etype=etype, bbox=bbox, section=block.section,
                    doc_id=normalized.doc_id,
                ))
        return elements

    @staticmethod
    def _reconstruct_text(block: dict) -> str:
        lines = []
        for line in block.get('lines', []):
            spans = [s.get('text', '') for s in line.get('spans', [])]
            t = ''.join(spans).strip()
            if t:
                lines.append(t)
        return '\n'.join(lines)

    @staticmethod
    def _classify_block(text: str, bbox: tuple, page_height: float) -> BlockType:
        if not bbox or len(bbox) < 4:
            return BlockType.TEXT
        if bbox[1] < page_height * 0.05:
            return BlockType.HEADER
        if bbox[3] > page_height * 0.95:
            return BlockType.FOOTER
        if text.strip().isupper() and len(text.split()) < 15:
            return BlockType.TITLE
        return BlockType.TEXT

    @staticmethod
    def _extract_tables(page_obj, page: NormalizedPage):
        try:
            import pdfplumber
            with pdfplumber.open(stream=page_obj.parent.to_bytes()) as pdf:
                pl_page = pdf.pages[page_obj.number]
                for table in pl_page.find_tables():
                    rows = table.extract()
                    if not rows or len(rows) < 2:
                        continue
                    page.blocks.append(NormalizedBlock(
                        type=BlockType.TABLE,
                        text=IngestWorkflow._rows_to_md(rows),
                        bbox=BoundingBox(*table.bbox),
                        page=page.page_number,
                    ))
        except Exception:
            pass

    @staticmethod
    def _rows_to_md(rows: list[list[str]]) -> str:
        rows = [[(c or '').strip() for c in r] for r in rows if r]
        if not rows:
            return ''
        h = rows[0]
        body = rows[1:] if len(rows) > 1 else []
        md = '| ' + ' | '.join(h) + ' |\n'
        md += '|' + '|'.join('---' for _ in h) + '|\n'
        for r in body:
            md += '| ' + ' | '.join((r + [''] * len(h))[: len(h)]) + ' |\n'
        return md

    def get_state(self) -> dict:
        return {
            'doc_id': self._doc_id,
            'filename': self._filename,
            'elements': self._elements,
            'tables': self._tables,
            'graph': self._graph,
            'vector_store': self.vector_store,
            'embedder': self.embedder,
            'geometry': self.geometry,
            'recurrence': self.recurrence,
            'frequency': self.frequency,
            'matrix': self.matrix,
            'template': self.template,
            'normalized': self._normalized,
        }
