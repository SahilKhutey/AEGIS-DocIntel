"""Tests for the document ingestion layer (loaders, OCR engine, and service routing)."""
import os
import sys
import io
import pytest
from pathlib import Path
from PIL import Image
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import Workbook
except ImportError:
    Workbook = None

try:
    import fitz
except ImportError:
    fitz = None

# Add amdi-os to sys.path
sys.path.insert(0, str(Path(__file__).parent.parent / "amdi-os"))

from src.core.document_object import DocumentObject, DocumentFormat
from src.ingestion import (
    IngestionService, OCREngine,
    PDFLoader, DOCXLoader, PPTXLoader, XLSXLoader, ImageLoader,
    FormatError, SizeLimitError, LoaderError
)

# ===== Helper Fixtures =====

@pytest.fixture
def minimal_pdf_bytes():
    if fitz is None:
        pytest.skip("PyMuPDF (fitz) is not installed")
    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((100, 100), "Hello PDF World")
    pdf_bytes = pdf.write()
    pdf.close()
    return pdf_bytes

@pytest.fixture
def minimal_docx_bytes():
    if DocxDocument is None:
        pytest.skip("python-docx is not installed")
    doc = DocxDocument()
    doc.add_paragraph("Hello DOCX Paragraph")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()

@pytest.fixture
def minimal_pptx_bytes():
    if Presentation is None:
        pytest.skip("python-pptx is not installed")
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX Title"
    buf = io.BytesIO()
    pres.save(buf)
    return buf.getvalue()

@pytest.fixture
def minimal_xlsx_bytes():
    if Workbook is None:
        pytest.skip("openpyxl is not installed")
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Hello XLSX Cell"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()

@pytest.fixture
def minimal_png_bytes():
    img = Image.new("RGB", (10, 10), color="blue")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

@pytest.fixture
def local_tmp_path():
    p = Path("tests/temp_test_dir")
    p.mkdir(parents=True, exist_ok=True)
    yield p
    # Clean up files and directory
    if p.exists():
        import shutil
        shutil.rmtree(p)


# ===== OCR Engine Tests =====

@pytest.mark.asyncio
async def test_ocr_empty_input():
    engine = OCREngine()
    result = await engine.recognize(b"")
    assert result == ""

@pytest.mark.asyncio
async def test_ocr_invalid_input():
    engine = OCREngine()
    result = await engine.recognize(b"not an image")
    assert result == ""

@pytest.mark.asyncio
async def test_ocr_png_image(minimal_png_bytes):
    engine = OCREngine()
    result = await engine.recognize(minimal_png_bytes)
    assert isinstance(result, str)

# ===== PDF Loader Tests =====

def test_pdf_validate_valid(minimal_pdf_bytes):
    loader = PDFLoader()
    assert loader.validate(minimal_pdf_bytes)
    assert loader.validate(b"%PDF-1.4\n...")
    assert loader.validate(b"%PDF")

def test_pdf_validate_invalid():
    loader = PDFLoader()
    assert not loader.validate(b"")
    assert not loader.validate(b"Not a PDF")
    assert not loader.validate(b"PK\x03\x04")

def test_pdf_validate_edge_cases():
    loader = PDFLoader()
    assert not loader.validate(b"%PD")
    assert not loader.validate(b"PDF-1.4")

@pytest.mark.asyncio
async def test_pdf_loader_minimal_pdf(minimal_pdf_bytes):
    loader = PDFLoader()
    doc = await loader.load(minimal_pdf_bytes, "test.pdf")
    assert doc.filename == "test.pdf"
    assert doc.format == DocumentFormat.PDF
    assert doc.size_bytes == len(minimal_pdf_bytes)
    assert doc.content_hash != ""
    assert doc.page_count == 1
    assert "page_count" in doc.metadata
    assert doc.metadata["scanned"] is False

@pytest.mark.asyncio
async def test_pdf_loader_size_limit(minimal_pdf_bytes):
    loader = PDFLoader(max_size_mb=0.000001)
    with pytest.raises(SizeLimitError):
        await loader.load(minimal_pdf_bytes, "big.pdf")

# ===== DOCX Loader Tests =====

def test_docx_validate(minimal_docx_bytes):
    loader = DOCXLoader()
    assert loader.validate(minimal_docx_bytes)
    assert not loader.validate(b"%PDF")

@pytest.mark.asyncio
async def test_docx_loader_minimal(minimal_docx_bytes):
    loader = DOCXLoader()
    doc = await loader.load(minimal_docx_bytes, "test.docx")
    assert doc.filename == "test.docx"
    assert doc.format == DocumentFormat.DOCX
    assert "Hello DOCX Paragraph" in doc.text_content
    assert doc.word_count > 0
    assert doc.metadata["paragraph_count"] > 0

# ===== PPTX Loader Tests =====

def test_pptx_validate(minimal_pptx_bytes):
    loader = PPTXLoader()
    assert loader.validate(minimal_pptx_bytes)

@pytest.mark.asyncio
async def test_pptx_loader_minimal(minimal_pptx_bytes):
    loader = PPTXLoader()
    doc = await loader.load(minimal_pptx_bytes, "test.pptx")
    assert doc.filename == "test.pptx"
    assert doc.format == DocumentFormat.PPTX
    assert "Hello PPTX Title" in doc.text_content
    assert doc.page_count == 1

# ===== XLSX Loader Tests =====

def test_xlsx_validate(minimal_xlsx_bytes):
    loader = XLSXLoader()
    assert loader.validate(minimal_xlsx_bytes)

@pytest.mark.asyncio
async def test_xlsx_loader_minimal(minimal_xlsx_bytes):
    loader = XLSXLoader()
    doc = await loader.load(minimal_xlsx_bytes, "test.xlsx")
    assert doc.filename == "test.xlsx"
    assert doc.format == DocumentFormat.XLSX
    assert "Hello XLSX Cell" in doc.text_content
    assert doc.page_count == 1

# ===== Image Loader Tests =====

def test_image_validate_formats(minimal_png_bytes):
    loader = ImageLoader()
    assert loader.validate(minimal_png_bytes)
    assert loader.validate(b"\xff\xd8\xff\xe0" + b"\x00" * 4)
    assert loader.validate(b"GIF89a" + b"\x00" * 2)
    assert loader.validate(b"BM" + b"\x00" * 6)
    assert loader.validate(b"II*\x00" + b"\x00" * 4)
    assert loader.validate(b"MM\x00*" + b"\x00" * 4)
    assert not loader.validate(b"%PDF")
    assert not loader.validate(b"")

@pytest.mark.asyncio
async def test_image_loader_minimal(minimal_png_bytes):
    loader = ImageLoader()
    doc = await loader.load(minimal_png_bytes, "test.png")
    assert doc.filename == "test.png"
    assert doc.format == DocumentFormat.IMAGE
    assert doc.page_count == 1
    assert doc.metadata["width"] == 10
    assert doc.metadata["height"] == 10

# ===== IngestionService Tests =====

def test_service_format_detection(minimal_pdf_bytes, minimal_docx_bytes, minimal_pptx_bytes, minimal_xlsx_bytes, minimal_png_bytes):
    service = IngestionService()
    assert service.detect_format(minimal_pdf_bytes, "test.pdf") == DocumentFormat.PDF
    assert service.detect_format(minimal_docx_bytes, "test.docx") == DocumentFormat.DOCX
    assert service.detect_format(minimal_pptx_bytes, "test.pptx") == DocumentFormat.PPTX
    assert service.detect_format(minimal_xlsx_bytes, "test.xlsx") == DocumentFormat.XLSX
    assert service.detect_format(minimal_png_bytes, "test.png") == DocumentFormat.IMAGE
    assert service.detect_format(b"Hello World", "test.txt") == DocumentFormat.TEXT
    assert service.detect_format(b"Hello World", "test.unknown") == DocumentFormat.UNKNOWN

def test_service_office_zip_detection_no_ext(minimal_docx_bytes, minimal_pptx_bytes, minimal_xlsx_bytes):
    service = IngestionService()
    assert service.detect_format(minimal_docx_bytes, "unknown") == DocumentFormat.DOCX
    assert service.detect_format(minimal_pptx_bytes, "unknown") == DocumentFormat.PPTX
    assert service.detect_format(minimal_xlsx_bytes, "unknown") == DocumentFormat.XLSX

def test_service_get_supported():
    service = IngestionService()
    formats = service.get_supported_formats()
    assert DocumentFormat.PDF in formats
    assert DocumentFormat.DOCX in formats
    assert DocumentFormat.PPTX in formats
    assert DocumentFormat.XLSX in formats
    assert DocumentFormat.IMAGE in formats

@pytest.mark.asyncio
async def test_service_ingest_bytes(minimal_pdf_bytes):
    service = IngestionService()
    doc = await service.ingest(minimal_pdf_bytes, "test.pdf")
    assert doc.format == DocumentFormat.PDF
    assert doc.filename == "test.pdf"

@pytest.mark.asyncio
async def test_service_ingest_file_path(local_tmp_path, minimal_pdf_bytes):
    service = IngestionService()
    pdf_file = local_tmp_path / "doc.pdf"
    pdf_file.write_bytes(minimal_pdf_bytes)
    
    doc = await service.ingest(pdf_file)
    assert doc.filename == "doc.pdf"
    assert doc.format == DocumentFormat.PDF

@pytest.mark.asyncio
async def test_service_batch_ingest(local_tmp_path, minimal_pdf_bytes, minimal_docx_bytes):
    service = IngestionService()
    f1 = local_tmp_path / "doc1.pdf"
    f1.write_bytes(minimal_pdf_bytes)
    f2 = local_tmp_path / "doc2.docx"
    f2.write_bytes(minimal_docx_bytes)
    
    docs = await service.ingest_batch([f1, f2])
    assert len(docs) == 2
    assert docs[0].filename == "doc1.pdf"
    assert docs[1].filename == "doc2.docx"

@pytest.mark.asyncio
async def test_invalid_formats_and_exceptions():
    service = IngestionService()
    with pytest.raises(LoaderError):
        await service.ingest(b"random content", "test.unknown")

    loader = PDFLoader()
    with pytest.raises(FormatError):
        await loader.load(b"not a pdf at all", "bad.pdf")
