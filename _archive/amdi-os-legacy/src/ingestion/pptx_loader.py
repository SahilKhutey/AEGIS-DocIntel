"""
AEGIS-AMDI-OS — PPTX Loader
=============================
Microsoft PowerPoint loader using python-pptx.
"""
from __future__ import annotations

import io
import logging
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.core.document_object import DocumentFormat, DocumentObject
from src.ingestion.base import BaseLoader, FormatError, SizeLimitError

logger = logging.getLogger(__name__)


class PPTXLoader(BaseLoader):
    """PPTX presentation loader."""

    FORMAT_NAME = "pptx"
    SUPPORTED_EXTENSIONS = {".pptx"}
    PPTX_MAGIC = b"PK\x03\x04"

    def __init__(self, max_size_mb: int = 100, **options):
        super().__init__(**options)
        self.max_size_mb = max_size_mb

    def validate(self, raw_bytes: bytes) -> bool:
        if not raw_bytes or len(raw_bytes) < 4:
            return False
        return raw_bytes[:4] == self.PPTX_MAGIC

    async def load(self, source, filename: str = "") -> DocumentObject:
        raw_bytes, name = self.read_source(source)
        if filename:
            name = filename
        if not name:
            name = "presentation.pptx"

        if not self.validate(raw_bytes):
            raise FormatError("Not a valid PPTX file")

        size_mb = len(raw_bytes) / (1024 * 1024)
        if size_mb > self.max_size_mb:
            raise SizeLimitError(f"PPTX too large: {size_mb:.1f}MB")

        metadata, text_parts = self._extract(raw_bytes)
        return DocumentObject(
            filename=name,
            format=DocumentFormat.PPTX,
            raw_bytes=raw_bytes,
            metadata=metadata,
            page_count=metadata.get("slide_count", 0),
            text_content="\n\n".join(text_parts),
        )

    def _extract(self, raw_bytes: bytes) -> tuple[dict[str, Any], list[str]]:
        """Extract metadata and text from PPTX."""
        metadata: dict[str, Any] = {}
        text_parts: list[str] = []
        try:
            pres = Presentation(io.BytesIO(raw_bytes))
            cp = pres.core_properties
            metadata["title"] = cp.title
            metadata["author"] = cp.author
            metadata["subject"] = cp.subject
            metadata["keywords"] = cp.keywords
            metadata["slide_count"] = len(pres.slides)
            metadata["slide_width"] = pres.slide_width
            metadata["slide_height"] = pres.slide_height

            image_count = 0
            chart_count = 0
            table_count = 0
            # Iterate slides
            for slide_idx, slide in enumerate(pres.slides, start=1):
                # Slide title
                if slide.shapes.title:
                    title_text = slide.shapes.title.text.strip()
                    if title_text:
                        text_parts.append(f"--- Slide {slide_idx}: {title_text} ---")
                # All shapes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text and text != (slide.shapes.title.text.strip() if slide.shapes.title else ""):
                                text_parts.append(text)
                    # Count shape types
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                    elif shape.has_chart:
                        chart_count += 1
                    elif shape.has_table:
                        table_count += 1
                        # Extract table data
                        tbl = shape.table
                        for row in tbl.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text:
                                text_parts.append(row_text)
            metadata["image_count"] = image_count
            metadata["chart_count"] = chart_count
            metadata["table_count"] = table_count
        except Exception as e:
            logger.warning(f"PPTX extraction failed: {e}")
        return metadata, text_parts
