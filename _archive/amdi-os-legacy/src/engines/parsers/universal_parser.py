"""
universal_parser.py
===================
AEGIS-AMDI-OS  |  Engine Layer  |  Document Parsing

Provides :class:`UniversalParser`, a format-agnostic document parser that
routes each :class:`~src.core.document_object.DocumentObject` to the
appropriate format-specific sub-parser and produces a
:class:`~src.core.normalized_document.NormalizedDocument`.

Supported formats
-----------------
* **PDF**  – PyMuPDF (fitz) for text/block extraction + pdfplumber for tables
* **DOCX** – python-docx
* **PPTX** – python-pptx
* **XLSX** – openpyxl
* **IMAGE** – stub (FIGURE block returned; OCR handled downstream)
* **TEXT / MARKDOWN / HTML** – plain-text split by double-newline

All third-party imports are wrapped in ``try/except ImportError`` so the
module loads even when optional packages are absent.
"""

from __future__ import annotations

import io
import logging
import os
import re
from pathlib import Path
from typing import List, Optional, Tuple

# ---------------------------------------------------------------------------
# Optional heavy dependencies – fail gracefully
# ---------------------------------------------------------------------------
try:
    import fitz  # PyMuPDF
    _FITZ_AVAILABLE = True
except ImportError:
    fitz = None  # type: ignore[assignment]
    _FITZ_AVAILABLE = False

try:
    import pdfplumber
    _PDFPLUMBER_AVAILABLE = True
except ImportError:
    pdfplumber = None  # type: ignore[assignment]
    _PDFPLUMBER_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    _DOCX_AVAILABLE = True
except ImportError:
    DocxDocument = None  # type: ignore[assignment]
    _DOCX_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Pt
    _PPTX_AVAILABLE = True
except ImportError:
    PptxPresentation = None  # type: ignore[assignment]
    Pt = None  # type: ignore[assignment]
    _PPTX_AVAILABLE = False

try:
    import openpyxl
    _OPENPYXL_AVAILABLE = True
except ImportError:
    openpyxl = None  # type: ignore[assignment]
    _OPENPYXL_AVAILABLE = False

try:
    import magic  # python-magic
    _MAGIC_AVAILABLE = True
except ImportError:
    magic = None  # type: ignore[assignment]
    _MAGIC_AVAILABLE = False

# ---------------------------------------------------------------------------
# Internal imports
# ---------------------------------------------------------------------------
from src.core.document_object import DocumentObject
from src.core.normalized_document import (
    NormalizedBlock,
    NormalizedDocument,
    NormalizedPage,
    BlockType,
)

# ---------------------------------------------------------------------------
# Logger
# ---------------------------------------------------------------------------
log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Magic-byte signatures for format detection
# ---------------------------------------------------------------------------
_MAGIC_BYTES: dict[bytes, str] = {
    b"%PDF": "pdf",
    b"PK\x03\x04": "zip",  # DOCX / XLSX / PPTX are ZIP-based
    b"\xff\xd8\xff": "image",
    b"\x89PNG": "image",
    b"GIF8": "image",
    b"BM": "image",
    b"II*\x00": "image",  # TIFF little-endian
    b"MM\x00*": "image",  # TIFF big-endian
}

_EXTENSION_MAP: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".tiff": "image",
    ".tif": "image",
    ".bmp": "image",
    ".txt": "text",
    ".md": "markdown",
    ".html": "html",
    ".htm": "html",
}


# ===========================================================================
# UniversalParser
# ===========================================================================
class UniversalParser:
    """
    Format-agnostic document parser.

    Routing logic
    -------------
    1. ``doc.format`` is inspected first (explicit caller hint).
    2. If absent or ``"auto"``, :meth:`_detect_format` is called.
    3. The resolved format is forwarded to the matching private method.

    Parameters
    ----------
    source_encoding : str
        Fallback text encoding used for plain-text formats (default ``utf-8``).
    max_pages : int | None
        Optional hard cap on pages parsed (useful for large PDFs in preview
        mode).  ``None`` means unlimited.
    """

    def __init__(
        self,
        source_encoding: str = "utf-8",
        max_pages: Optional[int] = None,
    ) -> None:
        self.source_encoding = source_encoding
        self.max_pages = max_pages
        log.debug(
            "UniversalParser initialised (encoding=%s, max_pages=%s)",
            source_encoding,
            max_pages,
        )

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    async def parse(self, doc: DocumentObject) -> NormalizedDocument:
        """
        Parse *doc* and return a :class:`NormalizedDocument`.

        Parameters
        ----------
        doc : DocumentObject
            The source document to parse.

        Returns
        -------
        NormalizedDocument
            Structured, page-segmented representation of *doc*.

        Raises
        ------
        ValueError
            If the document format cannot be determined.
        RuntimeError
            If the required parsing library is unavailable.
        """
        fmt = getattr(doc, "format", None) or "auto"
        if fmt == "auto":
            fmt = self._detect_format(doc)

        fmt = fmt.lower().strip()
        log.info("Parsing document id=%s as format=%s", getattr(doc, "doc_id", getattr(doc, "id", "?")), fmt)

        dispatch = {
            "pdf": self._parse_pdf,
            "docx": self._parse_docx,
            "pptx": self._parse_pptx,
            "xlsx": self._parse_xlsx,
            "image": self._parse_image,
            "text": self._parse_text,
            "txt": self._parse_text,
            "markdown": self._parse_text,
            "md": self._parse_text,
            "html": self._parse_text,
            "htm": self._parse_text,
        }

        handler = dispatch.get(fmt)
        if handler is None:
            log.warning("Unknown format '%s', falling back to plain-text parser.", fmt)
            handler = self._parse_text

        pages: List[NormalizedPage] = await handler(doc)

        # Apply max_pages cap
        if self.max_pages is not None:
            pages = pages[: self.max_pages]

        return NormalizedDocument(
            source_id=getattr(doc, "doc_id", getattr(doc, "id", None)),
            source_path=str(getattr(doc, "raw_path", getattr(doc, "path", "")) or ""),
            format=fmt,
            pages=pages,
            metadata=getattr(doc, "metadata", {}) or {},
        )

    # ------------------------------------------------------------------
    # Format detection
    # ------------------------------------------------------------------

    def _detect_format(self, doc: DocumentObject) -> str:
        """
        Detect the document format using magic bytes then file extension.

        Parameters
        ----------
        doc : DocumentObject
            Must expose either ``doc.content`` (``bytes``) or ``doc.path``
            (``str`` / ``Path``).

        Returns
        -------
        str
            Lower-cased format token, e.g. ``"pdf"``, ``"docx"``.
        """
        raw: Optional[bytes] = getattr(doc, "content", None)

        # --- magic-bytes detection ---
        if raw:
            header = raw[:8]
            for sig, fmt in _MAGIC_BYTES.items():
                if header.startswith(sig):
                    if fmt == "zip":
                        # Disambiguate DOCX / PPTX / XLSX by filename
                        ext = Path(getattr(doc, "path", "") or "").suffix.lower()
                        return _EXTENSION_MAP.get(ext, "docx")
                    return fmt

            # Try python-magic for deeper inspection
            if _MAGIC_AVAILABLE:
                try:
                    mime = magic.from_buffer(raw[:2048], mime=True)
                    _MIME_MAP = {
                        "application/pdf": "pdf",
                        "image/png": "image",
                        "image/jpeg": "image",
                        "image/tiff": "image",
                        "image/gif": "image",
                        "image/bmp": "image",
                        "text/plain": "text",
                        "text/html": "html",
                        "text/markdown": "markdown",
                    }
                    if mime in _MIME_MAP:
                        return _MIME_MAP[mime]
                except Exception as exc:  # pragma: no cover
                    log.debug("python-magic detection failed: %s", exc)

        # --- extension fallback ---
        path = Path(getattr(doc, "path", "") or "")
        ext = path.suffix.lower()
        fmt = _EXTENSION_MAP.get(ext)
        if fmt:
            return fmt

        log.warning(
            "Could not detect format for document '%s'; defaulting to 'text'.",
            getattr(doc, "path", "unknown"),
        )
        return "text"

    # ------------------------------------------------------------------
    # PDF parsing
    # ------------------------------------------------------------------

    async def _parse_pdf(self, doc: DocumentObject) -> List[NormalizedPage]:
        """
        Parse a PDF document using PyMuPDF (fitz).

        Block classification heuristics
        --------------------------------
        * **HEADER** – blocks whose top edge is in the upper 5 % of the page.
        * **FOOTER** – blocks whose bottom edge is in the lower 5 % of the page.
        * **TITLE**  – short (≤ 10 words), ALL-CAPS, or large-font text at the
          start of the page.
        * **FIGURE** – image blocks (``type == 1`` in PyMuPDF dict).
        * **TEXT**   – everything else.

        Tables are extracted separately with pdfplumber and injected as
        ``BlockType.TABLE`` blocks in markdown format.

        Parameters
        ----------
        doc : DocumentObject
            Source document.

        Returns
        -------
        list[NormalizedPage]
        """
        if not _FITZ_AVAILABLE:
            log.error(
                "PyMuPDF (fitz) is not installed. "
                "Install it with: pip install pymupdf"
            )
            return self._fallback_text_pages(doc)

        raw = self._get_raw(doc)
        pages: List[NormalizedPage] = []

        # Pre-extract tables via pdfplumber (keyed by page index)
        plumber_tables: dict[int, List[str]] = {}
        if _PDFPLUMBER_AVAILABLE and raw:
            try:
                with pdfplumber.open(io.BytesIO(raw)) as pdf_pl:
                    for pi, pg in enumerate(pdf_pl.pages):
                        tbls = pg.extract_tables() or []
                        md_tables = [
                            self._rows_to_markdown(tbl)
                            for tbl in tbls
                            if tbl
                        ]
                        if md_tables:
                            plumber_tables[pi] = md_tables
            except Exception as exc:
                log.warning("pdfplumber table extraction failed: %s", exc)

        try:
            if raw:
                pdf = fitz.open(stream=raw, filetype="pdf")
            else:
                path = str(getattr(doc, "path", ""))
                pdf = fitz.open(path)
        except Exception as exc:
            log.error("fitz failed to open document: %s", exc)
            return self._fallback_text_pages(doc)

        try:
            for page_idx, page in enumerate(pdf):
                page_rect = page.rect
                page_h = page_rect.height or 1.0
                page_w = page_rect.width or 1.0
                header_thresh = page_h * 0.05
                footer_thresh = page_h * 0.95

                blocks: List[NormalizedBlock] = []

                # Extract text blocks
                try:
                    raw_blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)[
                        "blocks"
                    ]
                except Exception as exc:
                    log.warning("fitz page %d text extraction failed: %s", page_idx, exc)
                    raw_blocks = []

                for blk in raw_blocks:
                    blk_type_id = blk.get("type", 0)
                    bbox = blk.get("bbox", (0, 0, 0, 0))  # (x0, y0, x1, y1)
                    x0, y0, x1, y1 = bbox

                    if blk_type_id == 1:
                        # Image block
                        blocks.append(
                            NormalizedBlock(
                                block_type=BlockType.FIGURE,
                                text="[IMAGE]",
                                bbox=(x0, y0, x1, y1),
                                page_index=page_idx,
                            )
                        )
                        continue

                    # Concatenate spans into plain text
                    text = ""
                    font_sizes: List[float] = []
                    for line in blk.get("lines", []):
                        for span in line.get("spans", []):
                            text += span.get("text", "")
                            font_sizes.append(span.get("size", 0.0))
                        text += "\n"
                    text = text.strip()
                    if not text:
                        continue

                    avg_font = (
                        sum(font_sizes) / len(font_sizes) if font_sizes else 0.0
                    )

                    # Classify block
                    block_type = self._classify_pdf_block(
                        text=text,
                        x0=x0,
                        y0=y0,
                        x1=x1,
                        y1=y1,
                        page_h=page_h,
                        page_w=page_w,
                        header_thresh=header_thresh,
                        footer_thresh=footer_thresh,
                        avg_font=avg_font,
                        block_index=len(blocks),
                    )

                    blocks.append(
                        NormalizedBlock(
                            block_type=block_type,
                            text=text,
                            bbox=(x0, y0, x1, y1),
                            page_index=page_idx,
                            font_size=avg_font,
                        )
                    )

                # Inject pdfplumber table blocks
                for md_table in plumber_tables.get(page_idx, []):
                    blocks.append(
                        NormalizedBlock(
                            block_type=BlockType.TABLE,
                            text=md_table,
                            bbox=(0, 0, page_w, page_h),
                            page_index=page_idx,
                        )
                    )

                pages.append(
                    NormalizedPage(
                        page_index=page_idx,
                        width=page_w,
                        height=page_h,
                        blocks=blocks,
                    )
                )
        finally:
            pdf.close()

        log.debug("PDF parsed: %d pages, %d total blocks",
                  len(pages), sum(len(p.blocks) for p in pages))
        return pages

    def _classify_pdf_block(
        self,
        *,
        text: str,
        x0: float,
        y0: float,
        x1: float,
        y1: float,
        page_h: float,
        page_w: float,
        header_thresh: float,
        footer_thresh: float,
        avg_font: float,
        block_index: int,
    ) -> BlockType:
        """
        Apply position and typographic heuristics to classify a PDF text block.

        Returns
        -------
        BlockType
        """
        # Header / footer by vertical position
        if y0 < header_thresh:
            return BlockType.HEADER
        if y1 > footer_thresh:
            return BlockType.FOOTER

        words = text.split()
        # TITLE heuristic: short, ALL_CAPS or large font at top of page
        is_short = len(words) <= 10
        is_all_caps = text.isupper() and len(text) > 2
        is_large_font = avg_font >= 14.0
        is_first_block = block_index == 0

        if is_short and (is_all_caps or is_large_font or is_first_block):
            return BlockType.TITLE

        return BlockType.TEXT

    # ------------------------------------------------------------------
    # DOCX parsing
    # ------------------------------------------------------------------

    async def _parse_docx(self, doc: DocumentObject) -> List[NormalizedPage]:
        """
        Parse a DOCX document using python-docx.

        Heading mapping
        ---------------
        * ``Heading 1`` → :attr:`BlockType.TITLE`
        * ``Heading 2..9`` → :attr:`BlockType.SUBTITLE`
        * ``Normal`` / body text → :attr:`BlockType.TEXT`
        * Tables → :attr:`BlockType.TABLE` (markdown)

        All content is placed on a single :class:`NormalizedPage`.

        Parameters
        ----------
        doc : DocumentObject

        Returns
        -------
        list[NormalizedPage]
        """
        if not _DOCX_AVAILABLE:
            log.error(
                "python-docx is not installed. "
                "Install it with: pip install python-docx"
            )
            return self._fallback_text_pages(doc)

        raw = self._get_raw(doc)
        try:
            if raw:
                docx_doc = DocxDocument(io.BytesIO(raw))
            else:
                path = str(getattr(doc, "path", ""))
                docx_doc = DocxDocument(path)
        except Exception as exc:
            log.error("python-docx failed to open document: %s", exc)
            return self._fallback_text_pages(doc)

        blocks: List[NormalizedBlock] = []
        y_cursor = 0.0
        BLOCK_H = 20.0  # synthetic height unit per block

        # Iterate body elements in document order
        from docx.oxml.ns import qn  # type: ignore[import-untyped]
        body = docx_doc.element.body

        for child in body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

            if tag == "p":
                # Paragraph
                from docx.text.paragraph import Paragraph  # type: ignore[import-untyped]
                para = Paragraph(child, docx_doc)
                text = para.text.strip()
                if not text:
                    y_cursor += BLOCK_H * 0.5
                    continue

                style_name: str = (
                    para.style.name if para.style and para.style.name else "Normal"
                )
                block_type = self._classify_docx_paragraph(style_name)

                blocks.append(
                    NormalizedBlock(
                        block_type=block_type,
                        text=text,
                        bbox=(0, y_cursor, 600, y_cursor + BLOCK_H),
                        page_index=0,
                    )
                )
                y_cursor += BLOCK_H

            elif tag == "tbl":
                # Table
                from docx.table import Table  # type: ignore[import-untyped]
                tbl = Table(child, docx_doc)
                rows: List[List[str]] = []
                for row in tbl.rows:
                    rows.append([cell.text.strip() for cell in row.cells])
                md = self._rows_to_markdown(rows)
                blocks.append(
                    NormalizedBlock(
                        block_type=BlockType.TABLE,
                        text=md,
                        bbox=(0, y_cursor, 600, y_cursor + BLOCK_H * len(rows)),
                        page_index=0,
                    )
                )
                y_cursor += BLOCK_H * len(rows)

        return [
            NormalizedPage(
                page_index=0,
                width=600.0,
                height=max(y_cursor, BLOCK_H),
                blocks=blocks,
            )
        ]

    @staticmethod
    def _classify_docx_paragraph(style_name: str) -> BlockType:
        """Map a python-docx style name to a :class:`BlockType`."""
        sl = style_name.lower()
        if "heading 1" in sl or "title" in sl:
            return BlockType.TITLE
        if re.match(r"heading\s+[2-9]", sl) or "subtitle" in sl:
            return BlockType.SUBTITLE
        return BlockType.TEXT

    # ------------------------------------------------------------------
    # PPTX parsing
    # ------------------------------------------------------------------

    async def _parse_pptx(self, doc: DocumentObject) -> List[NormalizedPage]:
        """
        Parse a PPTX presentation using python-pptx.

        Each slide becomes one :class:`NormalizedPage`.
        The title placeholder is mapped to :attr:`BlockType.TITLE`;
        all other text frames become :attr:`BlockType.TEXT`.

        Parameters
        ----------
        doc : DocumentObject

        Returns
        -------
        list[NormalizedPage]
        """
        if not _PPTX_AVAILABLE:
            log.error(
                "python-pptx is not installed. "
                "Install it with: pip install python-pptx"
            )
            return self._fallback_text_pages(doc)

        raw = self._get_raw(doc)
        try:
            if raw:
                prs = PptxPresentation(io.BytesIO(raw))
            else:
                path = str(getattr(doc, "path", ""))
                prs = PptxPresentation(path)
        except Exception as exc:
            log.error("python-pptx failed to open document: %s", exc)
            return self._fallback_text_pages(doc)

        pages: List[NormalizedPage] = []
        slide_w = prs.slide_width or 9144000  # EMU
        slide_h = prs.slide_height or 6858000

        for slide_idx, slide in enumerate(prs.slides):
            blocks: List[NormalizedBlock] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = "\n".join(
                    para.text for para in shape.text_frame.paragraphs
                ).strip()
                if not text:
                    continue

                left = shape.left or 0
                top = shape.top or 0
                width = shape.width or 0
                height = shape.height or 0

                # Determine block type
                is_title = (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
                block_type = BlockType.TITLE if is_title else BlockType.TEXT

                # Normalise EMU → points (1 pt = 12700 EMU)
                EMU = 12700
                blocks.append(
                    NormalizedBlock(
                        block_type=block_type,
                        text=text,
                        bbox=(
                            left / EMU,
                            top / EMU,
                            (left + width) / EMU,
                            (top + height) / EMU,
                        ),
                        page_index=slide_idx,
                    )
                )

            pages.append(
                NormalizedPage(
                    page_index=slide_idx,
                    width=slide_w / 12700,
                    height=slide_h / 12700,
                    blocks=blocks,
                )
            )

        log.debug("PPTX parsed: %d slides", len(pages))
        return pages

    # ------------------------------------------------------------------
    # XLSX parsing
    # ------------------------------------------------------------------

    async def _parse_xlsx(self, doc: DocumentObject) -> List[NormalizedPage]:
        """
        Parse an XLSX workbook using openpyxl.

        Each worksheet becomes one :class:`NormalizedPage` with a single
        :attr:`BlockType.TABLE` block containing the data in markdown format.

        Parameters
        ----------
        doc : DocumentObject

        Returns
        -------
        list[NormalizedPage]
        """
        if not _OPENPYXL_AVAILABLE:
            log.error(
                "openpyxl is not installed. "
                "Install it with: pip install openpyxl"
            )
            return self._fallback_text_pages(doc)

        raw = self._get_raw(doc)
        try:
            if raw:
                wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            else:
                path = str(getattr(doc, "path", ""))
                wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        except Exception as exc:
            log.error("openpyxl failed to open workbook: %s", exc)
            return self._fallback_text_pages(doc)

        pages: List[NormalizedPage] = []
        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            rows: List[List[str]] = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(cell) if cell is not None else "" for cell in row])

            md = self._rows_to_markdown(rows) if rows else "*Empty sheet*"
            # Synthetic page dimensions
            col_count = max((len(r) for r in rows), default=1)
            row_count = len(rows)
            page_w = max(col_count * 80.0, 400.0)
            page_h = max(row_count * 20.0, 100.0)

            pages.append(
                NormalizedPage(
                    page_index=sheet_idx,
                    width=page_w,
                    height=page_h,
                    blocks=[
                        NormalizedBlock(
                            block_type=BlockType.TABLE,
                            text=md,
                            bbox=(0, 0, page_w, page_h),
                            page_index=sheet_idx,
                            metadata={"sheet_name": sheet_name},
                        )
                    ],
                )
            )

        wb.close()
        log.debug("XLSX parsed: %d sheets", len(pages))
        return pages

    # ------------------------------------------------------------------
    # IMAGE parsing (stub)
    # ------------------------------------------------------------------

    async def _parse_image(self, doc: DocumentObject) -> List[NormalizedPage]:
        """
        Return a single-page document with one FIGURE block.

        Actual text extraction is deferred to the downstream OCR engine
        (:class:`~src.engines.ocr.universal_ocr.UniversalOCR`).

        Parameters
        ----------
        doc : DocumentObject

        Returns
        -------
        list[NormalizedPage]
        """
        raw = self._get_raw(doc)
        width, height = 800.0, 600.0

        # Attempt to read actual dimensions if Pillow is available
        try:
            from PIL import Image as PILImage  # type: ignore[import-untyped]
            img = PILImage.open(io.BytesIO(raw)) if raw else None
            if img:
                width, height = float(img.width), float(img.height)
        except Exception:
            pass

        log.debug("IMAGE stub parser: returning single FIGURE page.")
        return [
            NormalizedPage(
                page_index=0,
                width=width,
                height=height,
                blocks=[
                    NormalizedBlock(
                        block_type=BlockType.FIGURE,
                        text="[IMAGE – pending OCR]",
                        bbox=(0, 0, width, height),
                        page_index=0,
                        metadata={"raw_bytes_len": len(raw) if raw else 0},
                    )
                ],
            )
        ]

    # ------------------------------------------------------------------
    # Plain text / Markdown / HTML parsing
    # ------------------------------------------------------------------

    async def _parse_text(self, doc: DocumentObject) -> List[NormalizedPage]:
        """
        Parse plain text, Markdown, or HTML documents.

        The full content is treated as a single page.  Double-newlines
        (paragraph breaks) delimit individual :attr:`BlockType.TEXT` blocks.

        Parameters
        ----------
        doc : DocumentObject

        Returns
        -------
        list[NormalizedPage]
        """
        raw = self._get_raw(doc)
        if raw:
            try:
                text = raw.decode(self.source_encoding, errors="replace")
            except Exception:
                text = raw.decode("latin-1", errors="replace")
        else:
            text = getattr(doc, "text", "") or ""

        # Strip HTML tags for HTML format
        fmt = (getattr(doc, "format", "") or "").lower()
        if fmt in ("html", "htm"):
            text = re.sub(r"<[^>]+>", " ", text)
            text = re.sub(r"&[a-z]+;", " ", text)

        paragraphs = re.split(r"\n{2,}", text)
        blocks: List[NormalizedBlock] = []
        y_cursor = 0.0
        LINE_H = 14.0
        PAGE_W = 600.0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # Detect block type based on markdown patterns
            btype = BlockType.TEXT
            if para.startswith('#'):
                btype = BlockType.TITLE
            else:
                lines = [l.strip() for l in para.splitlines()]
                if len(lines) >= 2 and any(l.startswith('|') for l in lines):
                    table_lines = [l for l in lines if l.startswith('|')]
                    if len(table_lines) >= 2:
                        btype = BlockType.TABLE

            line_count = para.count("\n") + 1
            block_h = LINE_H * line_count
            blocks.append(
                NormalizedBlock(
                    block_type=btype,
                    text=para,
                    bbox=(0, y_cursor, PAGE_W, y_cursor + block_h),
                    page_index=0,
                )
            )
            y_cursor += block_h + LINE_H  # inter-paragraph gap

        if not blocks:
            blocks.append(
                NormalizedBlock(
                    block_type=BlockType.TEXT,
                    text="(empty document)",
                    bbox=(0, 0, PAGE_W, LINE_H),
                    page_index=0,
                )
            )

        return [
            NormalizedPage(
                page_index=0,
                width=PAGE_W,
                height=max(y_cursor, LINE_H),
                blocks=blocks,
            )
        ]

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _rows_to_markdown(rows: List[List[str]]) -> str:
        """
        Convert a 2-D list of strings to a GitHub-Flavored Markdown table.

        The first row is treated as the header row.  All subsequent rows are
        data rows.  ``None`` values and whitespace are normalised to empty
        strings.

        Parameters
        ----------
        rows : list[list[str]]
            2-D grid where ``rows[0]`` is the header.

        Returns
        -------
        str
            Markdown table string, e.g.::

                | col1 | col2 |
                |---|---|
                | val1 | val2 |

        Examples
        --------
        >>> UniversalParser._rows_to_markdown([["A", "B"], ["1", "2"]])
        '| A | B |\\n|---|---|\\n| 1 | 2 |'
        """
        if not rows:
            return ""

        def _clean(cell: object) -> str:
            return str(cell).replace("|", "\\|").strip() if cell is not None else ""

        def _row_to_md(cells: List[str]) -> str:
            return "| " + " | ".join(_clean(c) for c in cells) + " |"

        # Normalise column count
        col_count = max(len(r) for r in rows)
        normalised = [r + [""] * (col_count - len(r)) for r in rows]

        lines = [_row_to_md(normalised[0])]
        lines.append("|" + "---|" * col_count)
        for row in normalised[1:]:
            lines.append(_row_to_md(row))

        return "\n".join(lines)

    def _get_raw(self, doc: DocumentObject) -> Optional[bytes]:
        """
        Return raw bytes from *doc*, reading from disk if necessary.

        Parameters
        ----------
        doc : DocumentObject

        Returns
        -------
        bytes | None
        """
        raw: Optional[bytes] = getattr(doc, "content", getattr(doc, "raw_bytes", None))
        if raw:
            return raw

        path_str: str = str(getattr(doc, "path", getattr(doc, "raw_path", "")) or "")
        if path_str and Path(path_str).is_file():
            try:
                return Path(path_str).read_bytes()
            except OSError as exc:
                log.error("Failed to read file '%s': %s", path_str, exc)
        return None

    def _fallback_text_pages(self, doc: DocumentObject) -> List[NormalizedPage]:
        """
        Last-resort text extraction when the primary parser fails.

        Parameters
        ----------
        doc : DocumentObject

        Returns
        -------
        list[NormalizedPage]
        """
        raw = self._get_raw(doc)
        text = ""
        if raw:
            try:
                text = raw.decode(self.source_encoding, errors="replace")
            except Exception:
                text = raw.decode("latin-1", errors="replace")
        text = text.strip() or "(no extractable text)"
        return [
            NormalizedPage(
                page_index=0,
                width=600.0,
                height=600.0,
                blocks=[
                    NormalizedBlock(
                        block_type=BlockType.TEXT,
                        text=text,
                        bbox=(0, 0, 600, 600),
                        page_index=0,
                    )
                ],
            )
        ]
